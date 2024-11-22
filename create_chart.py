from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Data: Parent-Child relationships
data = pd.read_excel("data.xlsx").to_dict(orient="records")
#print(data)

# Create a PowerPoint presentation
prs = Presentation()

# Add a slide with a title layout
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Add a title
title = slide.shapes.title
title.text = "Dependency Tree - Organizational Chart"

# Add a textbox for each level
text_box = slide.shapes.add_textbox(100, 100, 800, 500)
text_frame = text_box.text_frame

# Build hierarchy
parents = {}
for row in data:
    parent = row["Parent"]
    child = row["Child"]
    if parent not in parents:
        parents[parent] = []
    parents[parent].append(child)

def add_hierarchy(parent, level=0):
    """Recursive function to add hierarchy to the text box"""
    paragraph = text_frame.add_paragraph()
    paragraph.text = ' ' * (level * 4) + str(parent)
    paragraph.font.size = Pt(16 - level)  # Adjust font size by level
    paragraph.alignment = PP_ALIGN.LEFT
    for child in parents.get(parent, []):
        add_hierarchy(child, level + 1)

# Start with the top-level parent
add_hierarchy(data[0]["Parent"])

# Save the presentation
prs.save("Dependency_Tree.pptx")
print("PowerPoint chart created: Dependency_Tree.pptx")
